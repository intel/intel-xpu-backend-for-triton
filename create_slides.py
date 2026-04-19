"""Generate PowerPoint slides about the tensor descriptor i32/i64 investigation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCD, 0xD6, 0xF4)
BLUE = RGBColor(0x89, 0xB4, 0xFA)
GREEN = RGBColor(0xA6, 0xE3, 0xA1)
RED = RGBColor(0xF3, 0x8B, 0xA8)
YELLOW = RGBColor(0xF9, 0xE2, 0xAF)
MAUVE = RGBColor(0xCB, 0xA6, 0xF7)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Consolas"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet(tf, text, font_size=16, color=LIGHT_GRAY, level=0, font_name="Consolas"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.level = level
    return p


# --- Slide 1: Title ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)
add_text_box(slide, 1, 1.5, 11, 1.5,
             "Tensor Descriptor vs Tensor of Pointers",
             font_size=36, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.2, 11, 1,
             "Pointer Arithmetic Bit Width Investigation",
             font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11, 1,
             "Intel XPU Backend for Triton",
             font_size=20, color=MAUVE, alignment=PP_ALIGN.CENTER)

# --- Slide 2: Problem Statement ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "Problem Statement", font_size=28, color=BLUE, bold=True)
tf = add_text_box(slide, 0.5, 1.2, 12, 5.5, "", font_size=16)
tf.paragraphs[0].text = ""

add_bullet(tf, "User writes: tl.make_tensor_descriptor() + desc.load()", color=GREEN)
add_bullet(tf, "rewrite_tensor_descriptor_to_pointer pass converts to tl.load(ptr, mask)", color=LIGHT_GRAY, level=1)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "User writes: tl.arange() + pointer arithmetic + tl.load(ptrs)", color=GREEN)
add_bullet(tf, "Direct tensor-of-pointer code, user controls bit widths", color=LIGHT_GRAY, level=1)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "Question: Is rewritten descriptor code always >= performance?", color=YELLOW, font_size=18)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "Finding: The rewrite pass forces ALL offset arithmetic to i64", color=RED, font_size=18)
add_bullet(tf, "MakeTensorDescOp shapes are Variadic<I32>, strides are Variadic<I64>", color=LIGHT_GRAY, level=1)
add_bullet(tf, "But rewrite castToI64() on shapes, indices, and ranges", color=LIGHT_GRAY, level=1)
add_bullet(tf, "Manual code uses tl.arange() (i32) * stride (i32) = i32", color=LIGHT_GRAY, level=1)

# --- Slide 3: Before - i64 Path ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "Before: All Arithmetic in i64", font_size=28, color=RED, bold=True)

code = """RewriteTensorDescriptorToPointer.cpp (before fix):

  getExpandedOffsetWithRange():
    indexRowType = i64                          // forced to i64
    ExtSIOp(make_range_i32) -> i64             // extend range
    AddIOp(splat_offset_i64, range_i64) -> i64 // i64 addition

  generatePtrFromOffsetRanges():
    MulIOp(offset_i64, stride_i64) -> i64      // i64 multiply
    AddPtrOp(ptr, offset_i64)                  // i64 offset

  generateMaskFromOffsetRanges():
    ConstantIntOp(0, i64)                      // i64 zero
    CmpIOp(sge, offset_i64, zero_i64)          // i64 comparison
    CmpIOp(slt, offset_i64, shape_i64)         // i64 comparison

  RewriteMakeTensorDesc:
    castToI64(shapes)                          // i32 -> i64 promotion"""

add_text_box(slide, 0.5, 1.2, 12, 6, code, font_size=14, color=LIGHT_GRAY)

# --- Slide 4: After - i32 Path ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "After: Keep i32, Extend Only at Stride Multiply",
             font_size=28, color=GREEN, bold=True)

code = """RewriteTensorDescriptorToPointer.cpp (after fix):

  getExpandedOffsetWithRange():
    indexI32RowType = i32                       // stay in i32
    AddIOp(splat_offset_i32, range_i32) -> i32  // i32 addition

  generatePtrFromOffsetRanges():
    ExtSIOp(offset_i32) -> i64                  // extend ONLY HERE
    MulIOp(offset_i64, stride_i64) -> i64       // i64 multiply (required)
    AddPtrOp(ptr, offset_i64)                   // i64 offset (required)

  generateMaskFromOffsetRanges():
    ConstantIntOp(0, i32)                       // i32 zero
    CmpIOp(sge, offset_i32, zero_i32)           // i32 comparison
    CmpIOp(slt, offset_i32, shape_i32)          // i32 comparison

  RewriteMakeTensorDesc:
    shapes passed through as-is (i32)           // no promotion"""

add_text_box(slide, 0.5, 1.2, 12, 6, code, font_size=14, color=LIGHT_GRAY)

# --- Slide 5: Register Pressure Impact ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "Register Pressure Impact", font_size=28, color=BLUE, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 12, 5.5, "", font_size=16)
tf.paragraphs[0].text = ""

add_bullet(tf, "Each i64 value occupies 2 GRF words vs 1 for i32", color=LIGHT_GRAY)
add_bullet(tf, "", font_size=8)

# Table-like content
add_bullet(tf, "Operation              Before (i64)    After (i32)", color=YELLOW, font_size=15)
add_bullet(tf, "-------------------------------------------------------", color=YELLOW, font_size=15)
add_bullet(tf, "Offset range (idx+range)  2 GRF/elem     1 GRF/elem", color=LIGHT_GRAY, font_size=15)
add_bullet(tf, "Shape bounds              2 GRF/elem     1 GRF/elem", color=LIGHT_GRAY, font_size=15)
add_bullet(tf, "Mask comparisons          i64 CmpIOp     i32 CmpIOp", color=LIGHT_GRAY, font_size=15)
add_bullet(tf, "Stride multiply           i64 x i64      ExtSI(i32) x i64", color=LIGHT_GRAY, font_size=15)
add_bullet(tf, "AddPtrOp offset           i64             i64 (unchanged)", color=LIGHT_GRAY, font_size=15)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "Net: ~50% register savings for intermediate offset computation", color=GREEN, font_size=18)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "Why this matters:", color=YELLOW)
add_bullet(tf, "Intel GPUs: higher register pressure -> more spills -> worse perf", color=LIGHT_GRAY, level=1)
add_bullet(tf, "Auto-GRF: spill_size > 1000 triggers recompile with 256-GRF mode", color=LIGHT_GRAY, level=1)
add_bullet(tf, "256-GRF halves thread occupancy -> fewer warps in flight", color=LIGHT_GRAY, level=1)

# --- Slide 6: IR Comparison ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "IR Comparison: Gather Example", font_size=28, color=BLUE, bold=True)

left_code = """Before (i64):

%range = tt.make_range : tensor<128xi32>
%ext = arith.extsi %range : i32 -> i64
%add = arith.addi %splat_i64, %ext : i64
%mul = arith.muli %add, %stride_i64 : i64

%zero_i64 = arith.constant 0 : i64
%cmp_lo = arith.cmpi sge, %add, %zero_i64
%shape_i64 = ... (cast from i32)
%cmp_hi = arith.cmpi slt, %add, %shape_i64"""

right_code = """After (i32):

%range = tt.make_range : tensor<128xi32>
%add = arith.addi %splat_i32, %range : i32
%ext = arith.extsi %add : i32 -> i64
%mul = arith.muli %ext, %stride_i64 : i64

%zero_i32 = arith.constant 0 : i32
%cmp_lo = arith.cmpi sge, %add, %zero_i32
%shape_i32 = ... (native i32)
%cmp_hi = arith.cmpi slt, %add, %shape_i32"""

add_text_box(slide, 0.3, 1.2, 6, 5.5, left_code, font_size=14, color=RED)
add_text_box(slide, 6.8, 1.2, 6, 5.5, right_code, font_size=14, color=GREEN)

# --- Slide 7: Validation ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "Validation", font_size=28, color=BLUE, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 12, 5.5, "", font_size=16)
tf.paragraphs[0].text = ""

add_bullet(tf, "Build: compile-triton.sh passed", color=GREEN, font_size=18)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "Python tests: 108/108 passed (test_block_tdesc.py)", color=GREEN, font_size=18)
add_bullet(tf, "Descriptor load/store with f32, f16, i8 and various shapes", color=LIGHT_GRAY, level=1)
add_bullet(tf, "Boundary padding (zeros in out-of-bounds region)", color=LIGHT_GRAY, level=1)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "Pre-commit checks: all passed (clang-format, etc.)", color=GREEN, font_size=18)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "Lit test updated: rewrite-tensor-descriptor-to-pointer.mlir", color=GREEN, font_size=18)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "Next steps:", color=YELLOW, font_size=18)
add_bullet(tf, "Performance benchmarking on target hardware", color=LIGHT_GRAY, level=1)
add_bullet(tf, "Compare GRF spill counts before/after", color=LIGHT_GRAY, level=1)
add_bullet(tf, "Matmul / attention kernel benchmarks", color=LIGHT_GRAY, level=1)

# --- Slide 8: Summary ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "Summary", font_size=28, color=BLUE, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 12, 5.5, "", font_size=16)
tf.paragraphs[0].text = ""

add_bullet(tf, "Investigation:", color=YELLOW, font_size=20)
add_bullet(tf, "Tensor descriptor rewrite forced all offset arithmetic to i64", color=LIGHT_GRAY, level=1)
add_bullet(tf, "Manual tensor-of-pointer code uses i32 (tl.arange -> i32)", color=LIGHT_GRAY, level=1)
add_bullet(tf, "~2x register pressure for intermediate offset computation", color=LIGHT_GRAY, level=1)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "Fix:", color=YELLOW, font_size=20)
add_bullet(tf, "Keep shapes (i32), indices (i32), ranges (i32) in native types", color=LIGHT_GRAY, level=1)
add_bullet(tf, "Extend to i64 ONLY at stride multiply (required for address math)", color=LIGHT_GRAY, level=1)
add_bullet(tf, "Mask comparisons stay in i32 (shapes are i32)", color=LIGHT_GRAY, level=1)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "File changed:", color=YELLOW, font_size=20)
add_bullet(tf, "RewriteTensorDescriptorToPointer.cpp (1 file, ~30 lines)", color=LIGHT_GRAY, level=1)

output_path = "/home/jovyan/intel-xpu-backend-for-triton/tensor_desc_investigation.pptx"
prs.save(output_path)
print(f"Slides saved to: {output_path}")
